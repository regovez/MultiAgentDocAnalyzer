import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def add_footer(slide):
    """Adds the GenAI disclaimer legend to the bottom of a slide."""
    # Positioning: 0.5 inches from left, 7.1 inches from top (standard slide height is 7.5)
    left = Inches(0.5)
    top = Inches(7.1)
    width = Inches(9)
    height = Inches(0.3)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "The PowerPoint was made by GenAI but based in the responses provided by the user."

    # Styling the footer to be subtle
    p.font.size = Pt(9)
    p.font.italic = True
    p.font.color.rgb = RGBColor(128, 128, 128)  # Grey color

def summarize_transcript_for_pptx(data, gpt4o):
    """
    Uses GPT-4o to convert raw Q&A into { 'Header': 'Summary' } format.
    """
    if not data:
        return {}

    prompt = f"""
    You are an Executive Editor. Transform this raw interview data into a professional summary table.

    Data: {json.dumps(data)}

    Instructions:
    1. For each Q&A pair, create a 'Key Header' (max 4 words) that captures the topic.
    2. For each answer, write a 'Distilled Summary' (max 20 words) that captures the core insight.
    3. Return ONLY a valid JSON object where the keys are the Headers and values are the Summaries.
    """

    response = gpt4o.invoke(prompt)
    try:
        clean_content = response.content.replace("```json", "").replace("```", "").strip()
        return json.loads(clean_content)
    except:
        return data


def generate_impactful_story(data, gpt4o):
    """
    Generates a narrative story for Slide 5 based on stakeholder instructions.
    """
    prompt = f"""
    Read all the answers received and create an impactful story for a single slide 
    to communicate the insights in a meaningful way. 

    Data: {json.dumps(data)}

    Instructions:
    1. Use only the facts/information shared without adding any additional info or data.
    2. Identify the key insights and client-relevant elements.
    3. The tone should be professional, persuasive, and cohesive.
    4. Keep the output concise enough to fit on a single PowerPoint slide (max 150 words).
    """

    response = gpt4o.invoke(prompt)
    return response.content.strip()


def create_executive_pptx(sub_id, submitter, transcript_json, gpt4o):
    # Ensure exports folder exists
    if not os.path.exists("exports"):
        os.makedirs("exports")

    prs = Presentation()
    data = json.loads(transcript_json)

    # Define key groups
    matrix_params = ["AI / GenAI", "Data", "Ecosystems", "Talent", "Methods", "Asset"]
    fixed_keys = ["Nature", "Persona", "USP", "Primary Market Challenge", "Avoidance"]

    # --- SLIDE 1: TITLE ---
    slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide_1.shapes.title.text = "Contextual Intelligence Report"
    slide_1.placeholders[1].text = f"Submitter: {submitter}"
    add_footer(slide_1)

    # --- SLIDE 2: STRATEGIC IMPACT MATRIX (Raw Data) ---
    slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide_2.shapes.title.text = "Strategic Impact Matrix"

    table_2 = add_styled_table(slide_2, rows=len(matrix_params) + 1, cols=2)
    table_2.cell(0, 0).text = "Dimension"
    table_2.cell(0, 1).text = "Impact Rating"

    for i, param in enumerate(matrix_params, start=1):
        table_2.cell(i, 0).text = param
        table_2.cell(i, 1).text = data.get(param, "Not Rated")

    add_footer(slide_2)
    # --- SLIDE 3: PROJECT FUNDAMENTALS (Raw Data) ---
    slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide_3.shapes.title.text = "Project Fundamentals"

    table_3 = add_styled_table(slide_3, rows=len(fixed_keys) + 1, cols=2)
    table_3.cell(0, 0).text = "Fundamental Parameter"
    table_3.cell(0, 1).text = "Detailed Response"

    for i, key in enumerate(fixed_keys, start=1):
        table_3.cell(i, 0).text = key
        raw_val = data.get(key, "N/A")
        # Clean multiselect string format
        table_3.cell(i, 1).text = raw_val.strip("[]").replace("'", "")

    add_footer(slide_3)
    # --- SLIDE 4: DISCOVERY INSIGHTS (AI Summarized) ---
    slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide_4.shapes.title.text = "Insights"

    # Filter data: only items NOT in the form keys (Questions 7-10)
    all_form_keys = matrix_params + fixed_keys
    chat_data = {q: a for q, a in data.items() if q not in all_form_keys}

    # Apply AI Summarization ONLY to this filtered chat data
    synthesized_chat = summarize_transcript_for_pptx(chat_data, gpt4o)

    table_4 = add_styled_table(slide_4, rows=len(synthesized_chat) + 1, cols=2)
    table_4.cell(0, 0).text = "Key Focus Area"
    table_4.cell(0, 1).text = "Executive Insight"

    for i, (header, summary) in enumerate(synthesized_chat.items(), start=1):
        table_4.cell(i, 0).text = header
        table_4.cell(i, 1).text = summary

    add_footer(slide_4)

    slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
    slide_5.shapes.title.text = "Strategic Narrative & Client Insights"

    # Generate the story using the new function
    story_text = generate_impactful_story(data, gpt4o)

    # Place text in the main content body
    body_shape = slide_5.placeholders[1]
    tf = body_shape.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = story_text
    p.font.size = Pt(14)  # Slightly smaller font to ensure the story fits

    # Add footer (if you kept that function)
    add_footer(slide_5)

    filename = f"exports/Strategic_Report_{sub_id}.pptx"
    prs.save(filename)
    return filename


def add_styled_table(slide, rows, cols):
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(0.8 * rows)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(6.5)
    return table